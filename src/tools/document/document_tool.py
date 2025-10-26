"""
Document Tool - Advanced file creation and reading capabilities

This tool provides comprehensive document handling including Markdown, Microsoft Word,
PowerPoint presentations, PDFs, and other common formats. It supports both creation
and reading operations with rich formatting capabilities.
"""

import os
from typing import Dict, Any, Optional, List
import sys

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '../..')))

from tools.base_tool import BaseTool


class DocumentTool(BaseTool):
    """
    Advanced document creation and reading tool.
    
    Supports multiple document formats including Markdown, Word documents,
    PowerPoint presentations, PDFs, and plain text files. Provides both
    creation and reading capabilities with formatting preservation.
    """
    
    SUPPORTED_FORMATS = {
        'markdown': ['.md', '.markdown'],
        'word': ['.docx'],
        'powerpoint': ['.pptx'],
        'pdf': ['.pdf'],
        'text': ['.txt'],
        'html': ['.html', '.htm']
    }
    
    def __init__(self, sandbox_path: str = "/tmp/graive_sandbox"):
        """
        Initialize the document tool.
        
        Args:
            sandbox_path: Path to the sandbox directory
        """
        self.sandbox_path = sandbox_path
        os.makedirs(sandbox_path, exist_ok=True)
        
    @property
    def name(self) -> str:
        return "document"
    
    @property
    def description(self) -> str:
        return (
            "Create and read documents in multiple formats including Markdown, "
            "Word (DOCX), PowerPoint (PPTX), PDF, HTML, and plain text"
        )
    
    @property
    def schema(self) -> Dict[str, Any]:
        return {
            "name": self.name,
            "description": self.description,
            "parameters": {
                "action": {
                    "type": "string",
                    "description": "Action to perform: create, read, convert",
                    "enum": ["create", "read", "convert"],
                    "required": True
                },
                "format": {
                    "type": "string",
                    "description": "Document format",
                    "enum": ["markdown", "word", "powerpoint", "pdf", "text", "html"],
                    "required": True
                },
                "filename": {
                    "type": "string",
                    "description": "Name of the file",
                    "required": True
                },
                "content": {
                    "type": "object",
                    "description": "Document content (structure varies by format)",
                    "required": False
                },
                "target_format": {
                    "type": "string",
                    "description": "Target format for conversion",
                    "required": False
                }
            }
        }
    
    def validate_parameters(self, parameters: Dict[str, Any]) -> bool:
        """Validate parameters for document operations."""
        required = ["action", "format", "filename"]
        if not all(key in parameters for key in required):
            return False
            
        action = parameters["action"]
        if action == "create" and "content" not in parameters:
            return False
            
        if action == "convert" and "target_format" not in parameters:
            return False
            
        return True
    
    def execute(self, parameters: Dict[str, Any]) -> Dict[str, Any]:
        """
        Execute document operation.
        
        Args:
            parameters: Operation parameters
            
        Returns:
            Operation result
        """
        action = parameters["action"]
        
        if action == "create":
            return self._create_document(parameters)
        elif action == "read":
            return self._read_document(parameters)
        elif action == "convert":
            return self._convert_document(parameters)
        else:
            return {"error": f"Unknown action: {action}", "success": False}
    
    def _create_document(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """Create a document in specified format."""
        doc_format = params["format"]
        filename = params["filename"]
        content = params["content"]
        
        if doc_format == "markdown":
            return self._create_markdown(filename, content)
        elif doc_format == "word":
            return self._create_word(filename, content)
        elif doc_format == "powerpoint":
            return self._create_powerpoint(filename, content)
        elif doc_format == "pdf":
            return self._create_pdf(filename, content)
        elif doc_format == "text":
            return self._create_text(filename, content)
        elif doc_format == "html":
            return self._create_html(filename, content)
        else:
            return {"error": f"Unsupported format: {doc_format}", "success": False}
    
    def _create_markdown(self, filename: str, content: Dict[str, Any]) -> Dict[str, Any]:
        """
        Create a Markdown document.
        
        Expected content structure:
        {
            "title": "Document Title",
            "sections": [
                {"heading": "Section 1", "level": 1, "content": "text"},
                {"heading": "Section 2", "level": 2, "content": "text"}
            ]
        }
        """
        try:
            filepath = os.path.join(self.sandbox_path, filename)
            
            md_content = []
            
            # Add title if provided
            if "title" in content:
                md_content.append(f"# {content['title']}\n")
            
            # Add sections
            if "sections" in content:
                for section in content["sections"]:
                    level = section.get("level", 2)
                    heading = section.get("heading", "")
                    text = section.get("content", "")
                    
                    md_content.append(f"{'#' * level} {heading}\n")
                    md_content.append(f"{text}\n")
            
            # Write to file
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write('\n'.join(md_content))
            
            return {
                "success": True,
                "message": f"Markdown document created: {filename}",
                "path": filepath
            }
            
        except Exception as e:
            return {"error": str(e), "success": False}
    
    def _create_word(self, filename: str, content: Dict[str, Any]) -> Dict[str, Any]:
        """
        Create a Word document (.docx).
        
        Expected content structure:
        {
            "title": "Document Title",
            "paragraphs": ["paragraph 1", "paragraph 2"],
            "headings": [{"text": "Heading", "level": 1}]
        }
        """
        try:
            from docx import Document
            from docx.shared import Pt, Inches
        except ImportError:
            return {
                "error": "python-docx not installed. Install with: pip install python-docx",
                "success": False
            }
        
        try:
            filepath = os.path.join(self.sandbox_path, filename)
            doc = Document()
            
            # Add title
            if "title" in content:
                doc.add_heading(content["title"], level=0)
            
            # Add headings and paragraphs
            if "sections" in content:
                for section in content["sections"]:
                    if "heading" in section:
                        doc.add_heading(section["heading"], level=section.get("level", 1))
                    if "content" in section:
                        doc.add_paragraph(section["content"])
            
            # Add simple paragraphs
            if "paragraphs" in content:
                for para in content["paragraphs"]:
                    doc.add_paragraph(para)
            
            doc.save(filepath)
            
            return {
                "success": True,
                "message": f"Word document created: {filename}",
                "path": filepath
            }
            
        except Exception as e:
            return {"error": str(e), "success": False}
    
    def _create_powerpoint(self, filename: str, content: Dict[str, Any]) -> Dict[str, Any]:
        """
        Create a PowerPoint presentation (.pptx).
        
        Expected content structure:
        {
            "title": "Presentation Title",
            "slides": [
                {
                    "type": "title",
                    "title": "Slide Title",
                    "subtitle": "Subtitle"
                },
                {
                    "type": "content",
                    "title": "Slide Title",
                    "points": ["Point 1", "Point 2"]
                }
            ]
        }
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return {
                "error": "python-pptx not installed. Install with: pip install python-pptx",
                "success": False
            }
        
        try:
            filepath = os.path.join(self.sandbox_path, filename)
            prs = Presentation()
            
            # Add slides
            if "slides" in content:
                for slide_content in content["slides"]:
                    slide_type = slide_content.get("type", "content")
                    
                    if slide_type == "title":
                        # Title slide
                        slide = prs.slides.add_slide(prs.slide_layouts[0])
                        slide.shapes.title.text = slide_content.get("title", "")
                        if slide.placeholders[1]:
                            slide.placeholders[1].text = slide_content.get("subtitle", "")
                    
                    elif slide_type == "content":
                        # Content slide
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        slide.shapes.title.text = slide_content.get("title", "")
                        
                        # Add bullet points
                        if "points" in slide_content:
                            text_frame = slide.placeholders[1].text_frame
                            text_frame.clear()
                            for point in slide_content["points"]:
                                p = text_frame.add_paragraph()
                                p.text = point
                                p.level = 0
            
            prs.save(filepath)
            
            return {
                "success": True,
                "message": f"PowerPoint presentation created: {filename}",
                "path": filepath
            }
            
        except Exception as e:
            return {"error": str(e), "success": False}
    
    def _create_pdf(self, filename: str, content: Dict[str, Any]) -> Dict[str, Any]:
        """Create a PDF document."""
        try:
            from fpdf import FPDF
        except ImportError:
            return {
                "error": "fpdf2 not installed. Install with: pip install fpdf2",
                "success": False
            }
        
        try:
            filepath = os.path.join(self.sandbox_path, filename)
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            
            # Add title
            if "title" in content:
                pdf.set_font("Arial", 'B', size=16)
                pdf.cell(0, 10, content["title"], ln=True)
                pdf.ln(5)
                pdf.set_font("Arial", size=12)
            
            # Add content
            if "sections" in content:
                for section in content["sections"]:
                    if "heading" in section:
                        pdf.set_font("Arial", 'B', size=14)
                        pdf.cell(0, 10, section["heading"], ln=True)
                        pdf.set_font("Arial", size=12)
                    if "content" in section:
                        pdf.multi_cell(0, 10, section["content"])
                        pdf.ln(5)
            
            pdf.output(filepath)
            
            return {
                "success": True,
                "message": f"PDF document created: {filename}",
                "path": filepath
            }
            
        except Exception as e:
            return {"error": str(e), "success": False}
    
    def _create_text(self, filename: str, content: Dict[str, Any]) -> Dict[str, Any]:
        """Create a plain text file."""
        try:
            filepath = os.path.join(self.sandbox_path, filename)
            
            text_content = []
            if "title" in content:
                text_content.append(content["title"])
                text_content.append("=" * len(content["title"]))
                text_content.append("")
            
            if "sections" in content:
                for section in content["sections"]:
                    if "heading" in section:
                        text_content.append(section["heading"])
                        text_content.append("-" * len(section["heading"]))
                    if "content" in section:
                        text_content.append(section["content"])
                        text_content.append("")
            
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write('\n'.join(text_content))
            
            return {
                "success": True,
                "message": f"Text file created: {filename}",
                "path": filepath
            }
            
        except Exception as e:
            return {"error": str(e), "success": False}
    
    def _create_html(self, filename: str, content: Dict[str, Any]) -> Dict[str, Any]:
        """Create an HTML document."""
        try:
            filepath = os.path.join(self.sandbox_path, filename)
            
            html_parts = ['<!DOCTYPE html>', '<html>', '<head>']
            
            if "title" in content:
                html_parts.append(f'<title>{content["title"]}</title>')
            
            html_parts.extend(['</head>', '<body>'])
            
            if "title" in content:
                html_parts.append(f'<h1>{content["title"]}</h1>')
            
            if "sections" in content:
                for section in content["sections"]:
                    level = section.get("level", 2)
                    if "heading" in section:
                        html_parts.append(f'<h{level}>{section["heading"]}</h{level}>')
                    if "content" in section:
                        html_parts.append(f'<p>{section["content"]}</p>')
            
            html_parts.extend(['</body>', '</html>'])
            
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write('\n'.join(html_parts))
            
            return {
                "success": True,
                "message": f"HTML document created: {filename}",
                "path": filepath
            }
            
        except Exception as e:
            return {"error": str(e), "success": False}
    
    def _read_document(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """Read a document and extract its content."""
        doc_format = params["format"]
        filename = params["filename"]
        filepath = os.path.join(self.sandbox_path, filename)
        
        if not os.path.exists(filepath):
            return {"error": f"File not found: {filename}", "success": False}
        
        if doc_format == "markdown" or doc_format == "text":
            return self._read_text_based(filepath)
        elif doc_format == "word":
            return self._read_word(filepath)
        elif doc_format == "powerpoint":
            return self._read_powerpoint(filepath)
        elif doc_format == "pdf":
            return self._read_pdf(filepath)
        elif doc_format == "html":
            return self._read_text_based(filepath)
        else:
            return {"error": f"Unsupported format: {doc_format}", "success": False}
    
    def _read_text_based(self, filepath: str) -> Dict[str, Any]:
        """Read text-based files (Markdown, TXT, HTML)."""
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                content = f.read()
            
            return {
                "success": True,
                "content": content,
                "path": filepath
            }
        except Exception as e:
            return {"error": str(e), "success": False}
    
    def _read_word(self, filepath: str) -> Dict[str, Any]:
        """Read Word document content."""
        try:
            from docx import Document
        except ImportError:
            return {
                "error": "python-docx not installed",
                "success": False
            }
        
        try:
            doc = Document(filepath)
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            
            return {
                "success": True,
                "content": {
                    "paragraphs": paragraphs,
                    "full_text": '\n'.join(paragraphs)
                },
                "path": filepath
            }
        except Exception as e:
            return {"error": str(e), "success": False}
    
    def _read_powerpoint(self, filepath: str) -> Dict[str, Any]:
        """Read PowerPoint presentation content."""
        try:
            from pptx import Presentation
        except ImportError:
            return {
                "error": "python-pptx not installed",
                "success": False
            }
        
        try:
            prs = Presentation(filepath)
            slides_content = []
            
            for slide in prs.slides:
                slide_data = {"shapes": []}
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_data["shapes"].append(shape.text)
                slides_content.append(slide_data)
            
            return {
                "success": True,
                "content": {
                    "slides": slides_content,
                    "slide_count": len(slides_content)
                },
                "path": filepath
            }
        except Exception as e:
            return {"error": str(e), "success": False}
    
    def _read_pdf(self, filepath: str) -> Dict[str, Any]:
        """Read PDF document content."""
        try:
            import PyPDF2
        except ImportError:
            return {
                "error": "PyPDF2 not installed. Install with: pip install PyPDF2",
                "success": False
            }
        
        try:
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                pages = []
                for page in reader.pages:
                    pages.append(page.extract_text())
            
            return {
                "success": True,
                "content": {
                    "pages": pages,
                    "page_count": len(pages),
                    "full_text": '\n'.join(pages)
                },
                "path": filepath
            }
        except Exception as e:
            return {"error": str(e), "success": False}
    
    def _convert_document(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """Convert document from one format to another."""
        # Read source document
        read_result = self._read_document(params)
        if not read_result.get("success"):
            return read_result
        
        # Create in target format
        target_format = params["target_format"]
        target_filename = params.get("target_filename", 
                                     params["filename"].rsplit('.', 1)[0] + 
                                     self.SUPPORTED_FORMATS.get(target_format, ['.txt'])[0])
        
        create_params = {
            "format": target_format,
            "filename": target_filename,
            "content": read_result["content"]
        }
        
        return self._create_document(create_params)
